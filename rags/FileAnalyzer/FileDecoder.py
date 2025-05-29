import os
import pathlib
import magic  # For MIME type detection
from llama_index.llms.gemini import Gemini
from llama_index.core.llms import ChatMessage, MessageRole

# --- Content Extractors ---
import pypdfium2 as pdfium
from docx import Document as DocxDocument
import pandas as pd
from pptx import Presentation
from bs4 import BeautifulSoup

from dotenv import load_dotenv

load_dotenv()

# Max characters to send to AI
MAX_CONTENT_CHARS = 50000
# Max rows to read from an Excel sheet
EXCEL_MAX_ROWS_TO_READ = 200


def extract_text_from_txt(file_path: pathlib.Path) -> str:
    """Extracts text from plain text files."""
    try:
        return file_path.read_text(encoding='utf-8')
    except UnicodeDecodeError:
        try:
            return file_path.read_text(encoding='latin-1')
        except UnicodeDecodeError:
            return file_path.read_text(encoding='utf-8', errors='ignore')
    except Exception as e:
        return f"[Error extracting TXT: {e}]"


def extract_text_from_pdf(file_path: pathlib.Path) -> str:
    """Extracts text from PDF files."""
    text = ""
    try:
        pdf = pdfium.PdfDocument(file_path)
        for i in range(len(pdf)):
            page = pdf.get_page(i)
            textpage = page.get_textpage()
            text += textpage.get_text_range() + "\n"
            textpage.close()
            page.close()
        pdf.close()
    except Exception as e:
        return f"[Error extracting PDF: {e}]"
    return text


def extract_text_from_docx(file_path: pathlib.Path) -> str:
    """Extracts text from DOCX files."""
    try:
        doc = DocxDocument(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        return f"[Error extracting DOCX: {e}]"


def extract_text_from_excel(file_path: pathlib.Path) -> str:
    """Extracts text from XLSX/XLS files."""
    try:
        # pandas uses 'openpyxl' for .xlsx and might try 'xlrd' for .xls.
        # Ensure 'openpyxl' is installed: pip install openpyxl
        # For .xls, if issues: pip install xlrd
        xls = pd.ExcelFile(file_path)
        text_parts = []

        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name, nrows=EXCEL_MAX_ROWS_TO_READ)

            header = f"Sheet: {sheet_name}\n"
            sheet_content_str = ""

            if df.empty:
                # Try to determine if the sheet was *actually* empty via the engine
                is_truly_empty = True  # Default assumption
                try:
                    sheet_engine_obj = xls.book[sheet_name]
                    if hasattr(sheet_engine_obj, 'max_row'):  # openpyxl (for .xlsx)
                        # max_row is 1-based index, or None if sheet is empty
                        if sheet_engine_obj.max_row is not None and sheet_engine_obj.max_row > 0:
                            is_truly_empty = False
                    elif hasattr(sheet_engine_obj, 'nrows'):  # xlrd (often for .xls)
                        # nrows is total number of rows
                        if sheet_engine_obj.nrows > 0:
                            is_truly_empty = False
                except Exception:
                    pass  # If inspection fails, rely on df.empty

                if is_truly_empty:
                    sheet_content_str = "(Sheet appears to be empty or contains no data cells)\n"
                else:
                    # df is empty, but underlying engine suggests rows exist. Pandas couldn't parse data.
                    sheet_content_str = f"(Pandas read no data from this sheet, though it might not be entirely empty. Attempted to read first {EXCEL_MAX_ROWS_TO_READ} rows.)\n"
            else:  # df is not empty
                sheet_content_str = df.to_string(index=False) + "\n"

                # Add truncation message if we read the maximum number of rows allowed
                if len(df) == EXCEL_MAX_ROWS_TO_READ:
                    truncation_note = f"[...displaying first {EXCEL_MAX_ROWS_TO_READ} rows. More rows might exist...]\n"
                    try:
                        sheet_engine_obj = xls.book[sheet_name]
                        if hasattr(sheet_engine_obj, 'max_row'):  # openpyxl
                            original_max_row = sheet_engine_obj.max_row
                            if original_max_row is not None and original_max_row > EXCEL_MAX_ROWS_TO_READ:
                                truncation_note = f"[...displaying first {EXCEL_MAX_ROWS_TO_READ} of approx. {original_max_row} data rows...]\n"
                        elif hasattr(sheet_engine_obj, 'nrows'):  # xlrd
                            original_nrows = sheet_engine_obj.nrows
                            if original_nrows > EXCEL_MAX_ROWS_TO_READ:
                                truncation_note = f"[...displaying first {EXCEL_MAX_ROWS_TO_READ} of {original_nrows} data rows...]\n"
                    except Exception as e_detail:
                        # Silently ignore if we can't get the exact count, the generic note is fine.
                        # print(f"(DEBUG: Could not get exact original row count for sheet '{sheet_name}': {e_detail})") # For debugging
                        pass
                    sheet_content_str += truncation_note

            text_parts.append(header + sheet_content_str)
        return "\n".join(text_parts)
    except ImportError as e:
        if 'xlrd' in str(e).lower():
            return "[Error extracting Excel: The 'xlrd' library is required for .xls files. Please install it: pip install xlrd]"
        if 'openpyxl' in str(e).lower():
            return "[Error extracting Excel: The 'openpyxl' library is required for .xlsx files. Please install it: pip install openpyxl]"
        return f"[Error extracting Excel: Missing a required library - {e}]"
    except Exception as e:
        err_str = str(e)
        if "File is not a zip file" in err_str and file_path.suffix.lower() == '.xlsx':  # xlsx are zip files
            return f"[Error extracting Excel: File '{file_path.name}' does not seem to be a valid XLSX (Zip) file. It might be corrupted or misnamed.]"
        if "Excel file format cannot be determined" in err_str:  # Pandas specific
            return f"[Error extracting Excel: Pandas could not determine the Excel file format for '{file_path.name}'. Ensure it's a valid .xls or .xlsx file and required libraries (xlrd/openpyxl) are installed.]"
        return f"[Error extracting Excel: {err_str}]"


def extract_text_from_pptx(file_path: pathlib.Path) -> str:
    """Extracts text from PPTX files."""
    try:
        prs = Presentation(file_path)
        text = ""
        for i, slide in enumerate(prs.slides):
            slide_text = f"--- Slide {i + 1} ---\n"
            has_text_in_slide = False
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text += shape.text + "\n"
                    has_text_in_slide = True
            if not has_text_in_slide:
                slide_text += "(No text found on this slide)\n"
            text += slide_text
        return text
    except Exception as e:
        return f"[Error extracting PPTX: {e}]"


def extract_text_from_html(file_path: pathlib.Path) -> str:
    """Extracts text content from HTML files."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            soup = BeautifulSoup(f, 'html.parser')
            text = ' '.join(soup.stripped_strings)
            return text
    except Exception as e:
        return f"[Error extracting HTML: {e}]"


def get_file_content(file_path_str: str) -> tuple[str | None, str | None]:
    """
    Detects file type and extracts its text content.
    Returns (content_string, error_message_string)
    """
    file_path = pathlib.Path(file_path_str)
    if not file_path.is_file():
        return None, f"Error: File not found at '{file_path_str}'"

    ext = file_path.suffix.lower()
    content: str | None = None
    error_msg: str | None = None
    mime_type: str | None = None

    try:
        # Mime type detection can be slow for large files, consider conditional use or timeout
        mime_type = magic.from_file(str(file_path), mime=True)
        print(f"Detected MIME type: {mime_type} (extension: {ext})")
    except Exception as e:
        print(f"Warning: Could not use python-magic: {e}. Relying primarily on extension.")

    # --- Priority 1: Extension-based for common structured documents & text ---
    if ext == ".docx":
        print("Processing as DOCX based on extension...")
        content = extract_text_from_docx(file_path)
    elif ext in [".xlsx", ".xls"]:
        print(f"Processing as Excel ({ext}) based on extension...")
        content = extract_text_from_excel(file_path)  # <-- FIXED FUNCTION CALLED HERE
    elif ext == ".pptx":
        print("Processing as PPTX based on extension...")
        content = extract_text_from_pptx(file_path)
    elif ext == ".pdf":
        print("Processing as PDF based on extension...")
        content = extract_text_from_pdf(file_path)
    elif ext in [".html", ".htm"]:
        print("Processing as HTML based on extension...")
        content = extract_text_from_html(file_path)
    elif ext in [".txt", ".py", ".json", ".csv", ".md", ".yaml", ".yml", ".log", ".srt", ".sub", ".xml", ".kml", ".gpx",
                 ".tsv"]:
        print(f"Processing as plain text ({ext}) based on extension...")
        content = extract_text_from_txt(file_path)
    elif ext in [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".webp", ".svg", ".ico"]:
        error_msg = f"File extension {ext} indicates an image. This script focuses on text content."

    # If content was extracted, check if it's an error message from the extractor
    if isinstance(content, str) and content.startswith("[Error extracting"):
        error_msg = content
        content = None  # Clear content as it's an error string

    # --- Priority 2: MIME-type based if extension didn't yield content or for other types ---
    if content is None and error_msg is None and mime_type:
        print(f"Attempting MIME-type based processing for: {mime_type}")
        if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" and ext != ".docx":
            content = extract_text_from_docx(file_path)
        elif mime_type in ["application/vnd.ms-excel",
                           "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"] and ext not in [".xlsx",
                                                                                                                ".xls"]:
            content = extract_text_from_excel(file_path)
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" and ext != ".pptx":
            content = extract_text_from_pptx(file_path)
        elif mime_type == "application/pdf" and ext != ".pdf":
            content = extract_text_from_pdf(file_path)
        elif (mime_type == "application/xhtml+xml" or mime_type == "text/html") and ext not in [".html", ".htm"]:
            content = extract_text_from_html(file_path)
        elif mime_type.startswith("text/"):  # General text types not caught by extension
            content = extract_text_from_txt(file_path)
        elif mime_type.startswith("image/"):
            error_msg = f"File is an image ({mime_type}). This script focuses on text content."
        elif mime_type.startswith("audio/") or mime_type.startswith("video/"):
            error_msg = f"File is an audio/video ({mime_type}). This script focuses on text content."
        elif mime_type == "application/zip":
            # If MIME is zip and extension didn't already handle it as an Office file (which are zips)
            if ext not in [".docx", ".xlsx", ".pptx", ".odp", ".ods", ".odt"]:  # Common office formats that are zips
                error_msg = (f"File MIME type is application/zip (extension: {ext}). "
                             "This is likely an archive. Text extraction from archives requires unpacking them first.")

        if isinstance(content, str) and content.startswith("[Error extracting"):
            error_msg = content
            content = None

    # --- Priority 3: Last resort - If no content and no specific error, try reading as text ---
    if content is None and error_msg is None:
        print(
            f"No specific handler for extension '{ext}' or MIME type '{mime_type}'. Attempting generic text extraction as a last resort...")
        # Avoid trying to read known binary/archive types as text if not already caught
        known_binary_or_archive_exts = [
            '.exe', '.dll', '.bin', '.dat', '.iso', '.img', '.zip', '.gz', '.tar',
            '.rar', '.7z', '.pkg', '.dmg', '.deb', '.rpm', '.jar', '.war', '.ear',
            '.gz', '.bz2', '.xz', '.apk', '.app', '.msi', '.so', '.o', '.a', '.lib',
            '.mp3', '.mp4', '.avi', '.mkv', '.mov', '.wav', '.flac', '.ogg', '.aac',
            '.woff', '.woff2', '.ttf', '.otf', '.eot'  # Font files
        ]  # Add more as needed
        if ext in known_binary_or_archive_exts or (mime_type and not mime_type.startswith("text/")):
            error_msg = f"File extension {ext} (MIME: {mime_type}) suggests a binary, archive, or non-text format not suitable for direct text reading."
        else:
            try:
                content = extract_text_from_txt(file_path)
                if isinstance(content, str) and content.startswith("[Error extracting"):
                    error_msg = content
                    content = None
                elif content is not None and not content.strip():
                    error_msg = f"File (ext: {ext}, MIME: {mime_type}) was read as text but resulted in empty or whitespace-only content."
                    content = None
            except Exception as e:
                error_msg = f"Unsupported file type (ext: {ext}, MIME: {mime_type}) or error during last resort text reading: {e}"

    if error_msg:
        return None, error_msg

    if content is None:
        return None, f"Could not extract text content from the file (ext: {ext}, MIME: {mime_type}). It might be a binary file or an unsupported format."

    if len(content) > MAX_CONTENT_CHARS:
        print(
            f"Warning: Content is very long ({len(content)} chars). Truncating to {MAX_CONTENT_CHARS} characters for AI.")
        content = content[:MAX_CONTENT_CHARS] + "\n[...content truncated...]"

    return content, None


def get_ai_description(text_content: str, model_name: str = "models/gemini-1.5-flash-latest") -> str:
    """ Sends text content to Gemini AI and returns its description. """
    api_key = os.getenv("GOOGLE_API_KEY")
    if not api_key:
        return "Error: GOOGLE_API_KEY environment variable not set."

    try:
        llm = Gemini(model_name=model_name, api_key=api_key)
        prompt = (
            "You are an AI assistant tasked with describing the contents of a file. "
            "Based on the following extracted text from the file, provide a concise summary "
            "and an overview of its content and purpose. \n"
            "If the content appears to be source code, describe what the code likely does, its language, and key functionalities. \n"
            "If it's structured data (like CSV, JSON, or Excel sheet data), describe the data's structure, an example of the kind of information it holds, and its potential use. \n"
            "If it's a document (like TXT, PDF, DOCX), summarize its main topic, key points, and intended audience or purpose. \n"
            "If the text is from a presentation (PPTX), summarize the overall topic and the key messages from the slides. \n"
            "If the text is from a webpage (HTML), describe the main content and purpose of the page. \n"
            "If the content is truncated, acknowledge that your summary is based on the available portion.\n"
            "If an error message is part of the extracted text (e.g., '[Error extracting...]'), mention that the extraction was partial or failed for that part.\n"
            "Be objective and stick to what can be inferred from the provided text.\n\n"
            "--- Extracted File Content ---\n"
            f"{text_content}"
            "\n\n--- End of File Content ---\n\n"
            "Detailed AI Description of File Contents:"
        )

        messages = [ChatMessage(role=MessageRole.USER, content=prompt)]
        response = llm.chat(messages)
        return response.message.content.strip()

    except Exception as e:
        return f"Error communicating with Gemini AI: {e}"


def main():
    file_path_str = input("Enter the full path to the file: ").strip()
    if not file_path_str:
        print("No file path entered. Exiting.")
        return

    try:
        # Remove leading/trailing quotes if present (common from drag-and-drop)
        if len(file_path_str) > 1 and file_path_str.startswith('"') and file_path_str.endswith('"'):
            file_path_str = file_path_str[1:-1]
        resolved_path = str(pathlib.Path(file_path_str).resolve(strict=True))
    except FileNotFoundError:
        print(f"Error: File not found at '{file_path_str}' (or after resolving potential quotes).")
        return
    except Exception as e:
        print(f"Error resolving path '{file_path_str}': {e}")
        return

    print(f"\nAttempting to process file: {resolved_path}")
    content, error_msg = get_file_content(resolved_path)

    if error_msg:
        print(f"\n--- Error during file processing ---")
        print(error_msg)
        # Decide if you want to send error messages to AI or just stop
        # For now, we stop. If you want to send, you'd do:
        # if content is None: content = f"[File processing error: {error_msg}]"
        # else: content += f"\n[File processing warning: {error_msg}]"
        # ... then proceed to AI description.
        return

    if not content:
        print(f"\n--- File Processing Result ---")
        print("No text content could be extracted, or the file is not suitable for text summarization.")
        return

    print(f"\n--- Extracted Content (first 500 chars) ---")
    print(content[:500] + ("..." if len(content) > 500 else ""))
    print("--- End of Extracted Content Sample ---")

    print("\nSending content to Gemini AI for description...")
    description = get_ai_description(content, model_name="models/gemini-2.0-flash")

    print(f"\n--- AI Description ---")
    print(description)


if __name__ == "__main__":
    # Ensure necessary libraries for pandas excel handling are hinted if missing
    try:
        import openpyxl
    except ImportError:
        print("Hint: For full .xlsx support, you might need to install 'openpyxl': pip install openpyxl")
    # try:
    #     import xlrd # xlrd is mainly for older .xls files
    # except ImportError:
    #     print("Hint: For older .xls file support, you might need to install 'xlrd': pip install xlrd")

    main()